# -*- coding: utf-8 -*-
"""PPT 개요(마크다운) → 실제 .pptx 변환. (회사 양식 템플릿 상속 지원)

Anthropic pptx 스킬의 '템플릿 기반 생성' 방식을 적용:
- 회사 양식 .pptx 가 있으면 그것을 베이스로 열어(테마·마스터·폰트·레이아웃 상속)
  기존 예시 슬라이드는 비우고, 양식의 레이아웃으로 새 슬라이드를 채운다.
- 양식이 없으면 기본 16:9 프레젠테이션으로 생성한다.

SYS_SCRIPT_PPT 개요 형식(### 슬라이드 N — 제목 + 불릿 + 발표자 노트)을 파싱한다.
"""
from __future__ import annotations

import io
import os
import re
from typing import Optional


def _clean(text: str) -> str:
    text = re.sub(r"\*\*|__|`|~~", "", text)
    text = re.sub(r"^[\-\*\+•·]\s*", "", text)
    return text.strip()


def _remove_all_slides(prs) -> None:
    """템플릿에 포함된 예시 슬라이드 제거(마스터·레이아웃은 보존).

    sldId 와 그 관계(rId)를 함께 끊는다. 도달 불가가 된 슬라이드 파트는
    저장 시 직렬화되지 않으므로 partname 충돌(Duplicate name)이 발생하지 않는다.
    """
    from pptx.oxml.ns import qn

    part = prs.part
    id_lst = prs.slides._sldIdLst
    for sid in list(id_lst):
        rid = sid.get(qn("r:id"))
        if rid and rid in part.rels:
            try:
                part.drop_rel(rid)
            except Exception:  # noqa: BLE001
                pass
        id_lst.remove(sid)


def _find_layout(prs, keywords, fallback_idx):
    for layout in prs.slide_layouts:
        nm = (layout.name or "").lower()
        if any(k in nm for k in keywords):
            return layout
    try:
        return prs.slide_layouts[fallback_idx]
    except Exception:  # noqa: BLE001
        return prs.slide_layouts[0]


def _body_placeholder(slide):
    """제목(idx 0) 이외의 첫 텍스트 플레이스홀더."""
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx != 0 and ph.has_text_frame:
                return ph
        except Exception:  # noqa: BLE001
            continue
    return None


def outline_to_pptx(md: str, deck_title: str = "강의 슬라이드",
                    template_path: Optional[str] = None) -> Optional[bytes]:
    """개요 마크다운 → .pptx 바이트. python-pptx 미설치 시 None."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except Exception:  # noqa: BLE001
        return None

    use_template = bool(template_path and os.path.isfile(template_path))
    if use_template:
        try:
            prs = Presentation(template_path)
            _remove_all_slides(prs)
        except Exception:  # noqa: BLE001
            prs = Presentation()
            use_template = False
    else:
        prs = Presentation()

    if not use_template:
        prs.slide_width = Pt(960)
        prs.slide_height = Pt(540)

    title_layout = _find_layout(prs, ["title slide", "제목 슬라이드", "표지", "cover"], 0)
    body_layout = _find_layout(prs, ["title and content", "제목 및 내용", "content", "내용", "본문"], 1)

    # ── 표지 ──
    s = prs.slides.add_slide(title_layout)
    if s.shapes.title is not None:
        s.shapes.title.text = deck_title
    sub = _body_placeholder(s)
    if sub is not None:
        sub.text_frame.text = "교수설계 가이드 에이전트 · Mayer 멀티미디어 원리 기반"

    # ── 슬라이드 블록 파싱 ──
    blocks = re.split(r"(?m)^\s*#{2,3}\s+", md or "")
    slide_blocks = [b for b in blocks if re.match(r"\s*슬라이드", b)]
    if not slide_blocks:
        slide_blocks = [b for b in blocks if b.strip()][:20]

    for block in slide_blocks:
        lines = block.splitlines()
        header = lines[0].strip() if lines else "슬라이드"
        title = re.sub(r"^슬라이드\s*\d+\s*[—\-:：]\s*", "", header).strip() or header

        body, notes, in_notes = [], [], False
        for ln in lines[1:]:
            t = ln.strip()
            if not t:
                continue
            if "발표자 노트" in t:
                in_notes = True
                after = re.split(r"[:：]", t, 1)
                if len(after) > 1 and after[1].strip():
                    notes.append(_clean(after[1]))
                continue
            if in_notes:
                notes.append(_clean(t))
            else:
                body.append(_clean(t))

        slide = prs.slides.add_slide(body_layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title[:120]
        bph = _body_placeholder(slide)
        if bph is not None and body:
            tf = bph.text_frame
            tf.clear()
            for i, line in enumerate(body[:12]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line[:200]
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = "\n".join(notes)
            except Exception:  # noqa: BLE001
                pass

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
