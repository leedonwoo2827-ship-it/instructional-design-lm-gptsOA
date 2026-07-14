# -*- coding: utf-8 -*-
"""디자인 슬라이드 빌더 — 슬라이드플랜(JSON) → 디자인된 .pptx 바이트.

애프터 덱(_guideline/after-...이미지·레이아웃정리.pptx)에서 실측한 네이비+앰버 디자인
시스템을 코드로 재현한다. 회사 템플릿을 베이스로 열어(테마·마스터·로고 상속) 예시
슬라이드를 비우고, 모든 슬라이드를 Blank 레이아웃에 도형으로 직접 배치한다.

타입: cover / section / photo / process / cards / compare / table / bullets
python-pptx 미설치 시 build_deck 은 None 을 돌려준다(그레이스풀).
"""
from __future__ import annotations

import io
import json
import os
import re
from typing import Callable, Dict, List, Optional

# ── 디자인 토큰(애프터 실측) ─────────────────────────────────────────────
FONT = "맑은 고딕"
SLIDE_W, SLIDE_H = 13.33, 7.5
MARGIN = 0.6
CONTENT_W = 12.13  # 전폭 콘텐츠

_HEX = {
    "navy": (0x1E, 0x27, 0x61),
    "navy_dk": (0x1B, 0x22, 0x52),
    "navy_dk2": (0x22, 0x2B, 0x63),
    "amber": (0xF2, 0xA9, 0x00),
    "chip": (0xEE, 0xF2, 0xFB),
    "teal": (0x2E, 0x7D, 0x8A),
    "white": (0xFF, 0xFF, 0xFF),
    "grey": (0x44, 0x49, 0x57),
    "sub": (0xCA, 0xDC, 0xFC),
}

# 프로세스 노드 / 카드 뱃지 색 로테이션
_NODE_ROT = ("navy_dk", "navy", "teal", "navy_dk2", "teal")
_BADGE_ROT = ("navy", "teal", "amber", "navy_dk")


def _rgb(name):
    from pptx.dml.color import RGBColor
    return RGBColor(*_HEX[name])


# ── 저수준 헬퍼 ───────────────────────────────────────────────────────────
def _remove_all_slides(prs) -> None:
    from pptx.oxml.ns import qn
    part = prs.part
    id_lst = prs.slides._sldIdLst
    for sid in list(id_lst):
        rid = sid.get(qn("r:id"))
        if rid and rid in part.rels:
            try:
                part.drop_rel(rid)
            except Exception:  # noqa: BLE001
                pass
        id_lst.remove(sid)


def _blank_layout(prs):
    for lay in prs.slide_layouts:
        nm = (lay.name or "").lower()
        if "blank" in nm or "빈" in (lay.name or ""):
            return lay
    return min(prs.slide_layouts, key=lambda l: len(l.placeholders))


def _strip_placeholders(slide) -> None:
    """레이아웃에서 상속된 빈 플레이스홀더 제거(마스터 로고·장식은 유지)."""
    for ph in list(slide.placeholders):
        try:
            ph._element.getparent().remove(ph._element)
        except Exception:  # noqa: BLE001
            pass


def _no_deco(shape):
    """도형 기본 그림자·테두리 제거."""
    try:
        shape.line.fill.background()
    except Exception:  # noqa: BLE001
        pass
    try:
        shape.shadow.inherit = False
    except Exception:  # noqa: BLE001
        pass


def _rrect(slide, x, y, w, h, fill=None, radius=None, line=None, line_w=None):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
    from pptx.oxml.ns import qn
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    _no_deco(shp)
    if fill is not None:
        shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(fill)
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = _rgb(line)
        shp.line.width = Pt(line_w or 1.0)
    if radius is not None:
        geom = shp._element.spPr.find(qn("a:prstGeom"))
        if geom is not None:
            av = geom.find(qn("a:avLst"))
            if av is None:
                av = geom.makeelement(qn("a:avLst"), {}); geom.append(av)
            for gd in list(av):
                av.remove(gd)
            gd = av.makeelement(qn("a:gd"), {"name": "adj", "fmla": f"val {int(radius*100000)}"})
            av.append(gd)
    return shp


def _oval(slide, x, y, d, fill):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    _no_deco(shp)
    shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(fill)
    return shp


def _diamond(slide, x, y, w, h, fill="amber"):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(x), Inches(y), Inches(w), Inches(h))
    _no_deco(shp)
    shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(fill)
    return shp


def _text(slide, x, y, w, h, anchor="t"):
    from pptx.util import Inches
    from pptx.enum.text import MSO_ANCHOR
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE,
                          "b": MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.TOP)
    return tf


def _run(p, text, size, color, bold=False, font=FONT):
    from pptx.util import Pt
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
    r.font.color.rgb = _rgb(color)
    return r


def _para(tf, first, align=None, space_after=4, line=1.05):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align:
        p.alignment = {"c": PP_ALIGN.CENTER, "l": PP_ALIGN.LEFT, "r": PP_ALIGN.RIGHT}[align]
    p.space_after = Pt(space_after); p.space_before = Pt(0)
    try:
        p.line_spacing = line
    except Exception:  # noqa: BLE001
        pass
    return p


# 제목·칩 세로 위치(v2 반영: 제목을 위에서 살짝 내리고 로고 피해 폭 축소)
TITLE_Y, TITLE_H, TITLE_W, TITLE_SZ = 0.7, 1.0, 10.6, 27
CHIP_Y = 1.82


# ── 구성 요소 ─────────────────────────────────────────────────────────────
def add_title(slide, text):
    tf = _text(slide, MARGIN, TITLE_Y, TITLE_W, TITLE_H, anchor="m")
    _run(_para(tf, True), text or "", TITLE_SZ, "navy", bold=True)


def add_chip(slide, text, x=MARGIN, y=CHIP_Y, w=CONTENT_W):
    """요약칩: 배경 라운드 + 앰버 점 + 네이비 굵은 한 문장. 반환=칩 하단 y."""
    if not text:
        return y
    h = 0.9
    _rrect(slide, x, y, w, h, fill="chip", radius=0.5)
    _oval(slide, x + 0.18, y + 0.36, 0.18, "amber")
    tf = _text(slide, x + 0.5, y, w - 0.7, h, anchor="m")
    _run(_para(tf, True, line=1.1), text, 15, "navy", bold=True)
    return y + h


def _hang(p, marL_in=0.23):
    """문단에 내어쓰기(hanging indent): ▸는 왼쪽으로 튀고 줄바꿈 글자는 텍스트 아래 정렬."""
    emu = str(int(marL_in * 914400))
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", emu)
    pPr.set("indent", "-" + emu)


def add_bullets(slide, bullets, x, y, w, h, size=15):
    if not bullets:
        return
    tf = _text(slide, x, y, w, h, anchor="t")
    for i, b in enumerate(bullets):
        p = _para(tf, i == 0, space_after=8, line=1.15)
        _hang(p)
        _run(p, "▸ ", size, "amber", bold=True)
        _run(p, str(b), size, "grey", bold=False)


def add_photo(slide, data, x=8.25, y=1.7, w=4.5, h=4.31):
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    try:
        pic = slide.shapes.add_picture(io.BytesIO(data), Inches(x), Inches(y),
                                       Inches(w), Inches(h))
    except Exception:  # noqa: BLE001
        return None
    # 종횡비 맞춰 센터 크롭(가능하면)
    try:
        from PIL import Image
        iw, ih = Image.open(io.BytesIO(data)).size
        tgt, src = w / h, iw / ih
        if src > tgt:
            c = (1 - tgt / src) / 2; pic.crop_left = c; pic.crop_right = c
        elif src < tgt:
            c = (1 - src / tgt) / 2; pic.crop_top = c; pic.crop_bottom = c
    except Exception:  # noqa: BLE001
        pass
    geom = pic._element.spPr.find(qn("a:prstGeom"))
    if geom is not None:
        geom.set("prst", "roundRect")
    return pic


def add_footer_key(slide, text, y=5.75):
    if not text:
        return
    h = 0.62
    _rrect(slide, MARGIN, y, CONTENT_W, h, fill="chip", radius=0.3)
    tf = _text(slide, MARGIN + 0.3, y, CONTENT_W - 0.5, h, anchor="m")
    p = _para(tf, True)
    _run(p, "핵심  ", 14, "amber", bold=True)
    _run(p, text, 14, "navy", bold=True)


def add_logo(slide, logo_path):
    from pptx.util import Inches
    if not logo_path or not os.path.isfile(logo_path):
        return
    try:
        slide.shapes.add_picture(logo_path, Inches(SLIDE_W - 1.85), Inches(0.28),
                                 height=Inches(0.42))
    except Exception:  # noqa: BLE001
        pass


# ── 타입별 렌더 ───────────────────────────────────────────────────────────
def render_cover(slide, s):
    _oval(slide, 9.6, -1.6, 5.6, "navy_dk")
    _oval(slide, 10.8, 3.6, 3.4, "navy_dk2")
    _rrect(slide, MARGIN, 2.15, 0.9, 0.9, fill="amber", radius=0.3)
    tf = _text(slide, MARGIN, 3.25, 10.8, 1.7, anchor="m")
    _run(_para(tf, True), s.get("title", "강의 슬라이드"), 40, "navy", bold=True)
    sub = s.get("chip") or s.get("subtitle") or ""
    if sub:
        tf2 = _text(slide, MARGIN, 4.75, 11.0, 0.7)
        _run(_para(tf2, True), sub, 18, "navy", bold=False)


def render_section(slide, s):
    _rrect(slide, MARGIN, 3.15, 0.16, 1.2, fill="amber", radius=0.5)
    tf = _text(slide, MARGIN + 0.4, 3.0, CONTENT_W - 0.4, 1.5, anchor="m")
    _run(_para(tf, True), s.get("title", ""), 34, "navy", bold=True)
    if s.get("chip"):
        tf2 = _text(slide, MARGIN + 0.4, 4.5, CONTENT_W - 0.4, 0.8)
        _run(_para(tf2, True), s["chip"], 16, "grey", bold=False)


# 사진 배치 프리셋(v2 실측): 슬라이드마다 순환해 크기·위치를 다양화.
# img=(L,T,W,H), tx/tw=텍스트(칩·본문) 열, below=이미지가 하단(본문은 위 전폭).
PHOTO_PRESETS = [
    {"img": (8.15, 2.5, 4.15, 4.15), "tx": 0.6, "tw": 7.3},   # 우측 정사각
    {"img": (0.5, 2.5, 4.15, 4.15), "tx": 4.95, "tw": 7.85},  # 좌측 정사각(텍스트 우)
    {"img": (10.0, 2.3, 2.83, 4.1), "tx": 0.6, "tw": 9.1},    # 우측 세로형
    {"img": (0.5, 2.3, 2.83, 4.1), "tx": 3.6, "tw": 9.2},     # 좌측 세로형(텍스트 우)
    {"img": (9.08, 2.7, 3.75, 3.75), "tx": 0.6, "tw": 8.2},   # 우측 중간
    # 하단 와이드(넓게) 프리셋 제거 — 전 슬라이드를 좌/우 사이드 사진으로 통일(요청).
]


def render_photo(slide, s, img):
    add_title(slide, s.get("title", ""))
    if not img:
        by = add_chip(slide, s.get("chip"))
        add_bullets(slide, s.get("bullets", []), MARGIN + 0.02, by + 0.15,
                    CONTENT_W, 6.95 - (by + 0.15))
        return
    p = PHOTO_PRESETS[s.get("_idx", 0) % len(PHOTO_PRESETS)]
    add_photo(slide, img, *p["img"])
    tx, tw = p["tx"], p["tw"]
    by = add_chip(slide, s.get("chip"), x=tx, w=tw)
    top = by + 0.15
    bottom = (p["img"][1] - 0.15) if p.get("below") else 6.95
    add_bullets(slide, s.get("bullets", []), tx + 0.02, top, tw - 0.04, max(1.0, bottom - top))


def render_bullets(slide, s):
    add_title(slide, s.get("title", ""))
    by = add_chip(slide, s.get("chip"))
    add_bullets(slide, s.get("bullets", []), MARGIN + 0.02, max(by + 0.2, 2.72),
                CONTENT_W, 6.9 - max(by + 0.2, 2.72))


def render_process(slide, s):
    add_title(slide, s.get("title", ""))
    add_chip(slide, s.get("chip"))
    items = s.get("items") or [{"label": b} for b in s.get("bullets", [])]
    items = items[:5] or [{"label": "항목"}]
    n = len(items)
    gap = 0.34
    node_h = 1.5
    node_w = (CONTENT_W - gap * (n - 1)) / n
    y = 3.15
    for i, it in enumerate(items):
        x = MARGIN + i * (node_w + gap)
        color = _NODE_ROT[i % len(_NODE_ROT)]
        _rrect(slide, x, y, node_w, node_h, fill=color, radius=0.12)
        tf = _text(slide, x + 0.1, y, node_w - 0.2, node_h, anchor="m")
        p = _para(tf, True, align="c")
        _run(p, it.get("label", ""), 15, "white", bold=True)
        desc = it.get("desc")
        if desc:
            tf2 = _text(slide, x + 0.05, y + node_h + 0.05, node_w - 0.1, 1.0, anchor="t")
            _run(_para(tf2, True, align="c", line=1.1), desc, 11, "grey", bold=False)
        if i < n - 1:
            _diamond(slide, x + node_w - 0.13, y + node_h / 2 - 0.16, 0.26, 0.32, "amber")
    add_footer_key(slide, s.get("key") or s.get("footer"))


def render_cards(slide, s):
    add_title(slide, s.get("title", ""))
    add_chip(slide, s.get("chip"))
    items = s.get("items") or [{"label": b} for b in s.get("bullets", [])]
    items = items[:4] or [{"label": "항목"}]
    n = len(items)
    gap = 0.34
    card_w = (CONTENT_W - gap * (n - 1)) / n
    y, card_h = 2.95, 2.7
    for i, it in enumerate(items):
        x = MARGIN + i * (card_w + gap)
        _rrect(slide, x, y, card_w, card_h, fill="white", radius=0.1,
               line="chip", line_w=1.5)
        badge = _BADGE_ROT[i % len(_BADGE_ROT)]
        _oval(slide, x + 0.28, y + 0.28, 0.55, badge)
        tfb = _text(slide, x + 0.28, y + 0.28, 0.55, 0.55, anchor="m")
        _run(_para(tfb, True, align="c"), str(i + 1), 16, "white", bold=True)
        tft = _text(slide, x + 0.28, y + 1.0, card_w - 0.56, 0.5)
        _run(_para(tft, True), it.get("label", ""), 16, "navy", bold=True)
        if it.get("desc"):
            tfd = _text(slide, x + 0.28, y + 1.55, card_w - 0.56, card_h - 1.7)
            _run(_para(tfd, True, line=1.15), it["desc"], 12.5, "grey", bold=False)
    add_footer_key(slide, s.get("key") or s.get("footer"))


def render_compare(slide, s):
    add_title(slide, s.get("title", ""))
    add_chip(slide, s.get("chip"))
    items = (s.get("items") or [])[:2]
    while len(items) < 2:
        items.append({"label": "", "desc": ""})
    y, h = 2.95, 3.6
    col_w = (CONTENT_W - 0.4) / 2
    for i, it in enumerate(items):
        x = MARGIN + i * (col_w + 0.4)
        _rrect(slide, x, y, col_w, h, fill="chip", radius=0.08)
        tft = _text(slide, x + 0.3, y + 0.25, col_w - 0.6, 0.6)
        _run(_para(tft, True), it.get("label", ""), 17, "navy", bold=True)
        lines = it.get("lines") or ([it.get("desc")] if it.get("desc") else [])
        if lines:
            add_bullets(slide, lines, x + 0.3, y + 1.0, col_w - 0.6, h - 1.2, size=13.5)


def render_table(slide, s):
    from pptx.util import Inches, Pt
    add_title(slide, s.get("title", ""))
    by = add_chip(slide, s.get("chip"))
    rows_data = s.get("rows")
    if not rows_data:
        return render_bullets(slide, s)
    headers = s.get("headers") or []
    y = max(by + 0.25, 2.8)
    ncol = len(headers) or max(len(r) for r in rows_data)
    nrow = len(rows_data) + (1 if headers else 0)
    tbl_shape = slide.shapes.add_table(nrow, ncol, Inches(MARGIN), Inches(y),
                                       Inches(CONTENT_W), Inches(min(0.5 * nrow, 4.0)))
    tbl = tbl_shape.table
    r0 = 0
    if headers:
        for c, htxt in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = str(htxt)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.bold = True; run.font.size = Pt(12); run.font.name = FONT
        r0 = 1
    for ri, row in enumerate(rows_data):
        for c in range(ncol):
            cell = tbl.cell(ri + r0, c)
            cell.text = str(row[c]) if c < len(row) else ""
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11); run.font.name = FONT


_RENDER = {
    "cover": lambda sl, s, img: render_cover(sl, s),
    "section": lambda sl, s, img: render_section(sl, s),
    "photo": render_photo,
    "process": lambda sl, s, img: render_process(sl, s),
    "cards": lambda sl, s, img: render_cards(sl, s),
    "compare": lambda sl, s, img: render_compare(sl, s),
    "table": lambda sl, s, img: render_table(sl, s),
    "bullets": lambda sl, s, img: render_bullets(sl, s),
}


def build_deck(plan: List[Dict], template_path: Optional[str] = None,
               images: Optional[Dict[int, bytes]] = None,
               deck_title: str = "강의 슬라이드",
               logo_path: Optional[str] = None) -> Optional[bytes]:
    """슬라이드플랜 → 디자인된 .pptx 바이트. python-pptx 미설치 시 None."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:  # noqa: BLE001
        return None

    images = images or {}
    use_tpl = bool(template_path and os.path.isfile(template_path))
    if use_tpl:
        try:
            prs = Presentation(template_path)
            _remove_all_slides(prs)
        except Exception:  # noqa: BLE001
            prs = Presentation(); use_tpl = False
    else:
        prs = Presentation()
    if not use_tpl:
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)

    layout = _blank_layout(prs)
    for i, s in enumerate(plan):
        slide = prs.slides.add_slide(layout)
        _strip_placeholders(slide)
        s["_idx"] = i  # 사진 배치 프리셋 순환용
        typ = (s.get("type") or "bullets").lower()
        fn = _RENDER.get(typ, _RENDER["bullets"])
        try:
            fn(slide, s, images.get(i))
        except Exception as e:  # noqa: BLE001
            print(f"[deck] slide {i} ({typ}) 렌더 오류: {e}", flush=True)
            try:
                render_bullets(slide, s)
            except Exception:  # noqa: BLE001
                pass
        add_logo(slide, logo_path)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── 아트디렉터(LLM) 패스 — 개요 md → 슬라이드플랜 JSON ─────────────────────
_ART_SYS = (
    "너는 강의 슬라이드 아트디렉터다. 주어진 슬라이드 개요를 슬라이드플랜 JSON 배열로만 변환한다. "
    "설명·머리말·코드펜스 없이 JSON 배열만 출력한다."
)

_ART_RULES = """개요의 '### 슬라이드 N' 블록 하나당 JSON 객체 하나를 순서대로 만든다(개수 유지).
각 객체 스키마:
{
  "type": "section|photo|process|cards|compare|table|bullets",
  "title": "슬라이드 제목",
  "chip": "요약 한 문장(핵심 메시지)",
  "bullets": ["불릿 3~5개"],
  "items": [{"label":"짧은 이름","desc":"한 줄 설명"}],
  "image_query": "photo형일 때만, 영어 검색어(예: 'classroom students learning')"
}
타입 선택 규칙(우선순위 순):
- "cover"는 만들지 마라(표지는 시스템이 자동으로 맨 앞에 추가한다).
- **"photo"를 최우선으로, 전체 슬라이드의 약 55~65%에 사용하라.** 사진으로 보완될 여지가 조금이라도 있으면 photo 를 선택하고 영어 image_query 를 넣는다(bullets 3~4개).
  · image_query 는 **주제·이론과 직접 관련된 구체적 영어 검색어**로(막연한 'education' 금지). 예) 행동주의→'Pavlov classical conditioning dog experiment', 인지주의→'human brain memory diagram', 구성주의→'students collaborative project based learning', 매체→'classroom multimedia projector lesson', 평가→'students taking written exam', 커뮤니케이션→'Shannon Weaver communication model', 플립러닝→'flipped classroom students laptop', 원격교육→'online video lecture student'. 강의실·학생·교사·실험·기기 장면은 거의 항상 사진 가능.
- 구성요소/단계/절차의 '순서'가 핵심인 것만: "process" (items[].label 2~5개).
- 3~4개의 병렬 개념·분류·유형: "cards" (items[].label + desc).
- 두 개념 대비/비교: "compare" (items 2개, 각 {"label","lines":[...]}).
- 표가 꼭 필요한 것(구성 개요·평가 유형 등): "table" (headers:[...], rows:[[...]]).
- 도입·구간 전환: "section".
- 위 어디에도 안 맞을 때만: "bullets".
공통: 모든 슬라이드에 chip(요약 한 문장). 같은 타입이 여러 장 연속되지 않게 섞되 photo 비중을 40~50%로 유지. process/cards엔 key(하단 핵심 한 줄)를 넣어도 좋다.
"""


def _extract_json_array(text: str):
    """JSON 배열 파싱. 잘린 응답도 완결된 {…} 객체만 골라 살려낸다."""
    if not text:
        return None
    t = re.sub(r"^```[a-zA-Z]*\s*|\s*```$", "", text.strip())
    i = t.find("[")
    if i == -1:
        return None
    j = t.rfind("]")
    if j > i:
        try:
            return json.loads(t[i:j + 1])
        except Exception:  # noqa: BLE001
            pass
    # 살리기: 배열 안의 완결된 최상위 { … } 객체들을 순서대로 파싱
    objs, depth, start, instr, esc = [], 0, None, False, False
    for k in range(i + 1, len(t)):
        c = t[k]
        if instr:
            if esc:
                esc = False
            elif c == "\\":
                esc = True
            elif c == '"':
                instr = False
            continue
        if c == '"':
            instr = True
        elif c == "{":
            if depth == 0:
                start = k
            depth += 1
        elif c == "}":
            if depth > 0:
                depth -= 1
                if depth == 0 and start is not None:
                    try:
                        objs.append(json.loads(t[start:k + 1]))
                    except Exception:  # noqa: BLE001
                        pass
                    start = None
    return objs or None


def _outline_blocks(outline_md: str) -> List[str]:
    """개요 md 를 '### 슬라이드 …' 단위 블록 리스트로 분할."""
    parts = re.split(r"(?m)^(?=#{2,3}\s*슬라이드)", outline_md or "")
    return [p.strip() for p in parts if re.match(r"#{2,3}\s*슬라이드", p.strip())]


def _block_to_bullets(block: str) -> Dict:
    """개요 블록 하나 → bullets 슬라이드 dict(아트디렉터 실패 시 폴백)."""
    lines = [l.strip() for l in block.splitlines() if l.strip()]
    title = re.sub(r"^#{2,3}\s*", "", lines[0]) if lines else "슬라이드"
    title = re.sub(r"^슬라이드\s*\d+\s*[—\-:：]\s*", "", title).strip()
    chip, bullets = "", []
    for ln in lines[1:]:
        m = re.match(r"^[-*+]?\s*\*{0,2}([^*:：]{1,20})\*{0,2}\s*[:：]\s*(.*)$", ln)
        if m and "핵심" in m.group(1):
            chip = m.group(2).strip()
        elif m and "레이아웃" in m.group(1):
            continue
        else:
            v = re.sub(r"^[-*+•·]\s*", "", ln)
            if v and not re.match(r"^\**본문", v):
                bullets.append(v)
    return {"type": "bullets", "title": title, "chip": chip, "bullets": bullets[:5]}


def _fallback_plan(outline_md: str, deck_title: str) -> List[Dict]:
    """개요 md 를 단순 파싱해 bullets 위주 플랜으로(아트디렉터 실패 시)."""
    plan = [{"type": "cover", "title": deck_title, "chip": ""}]
    plan += [_block_to_bullets(b) for b in _outline_blocks(outline_md)]
    return plan


def _normalize_slide(s: Dict) -> Optional[Dict]:
    if not isinstance(s, dict):
        return None
    s.setdefault("type", "bullets")
    s.setdefault("title", "")
    s.setdefault("bullets", [])
    if s["type"] == "cover":   # LLM이 표지를 만들면 내용형으로 강등
        s["type"] = "bullets"
    return s


_CHUNK = 14  # 청크당 슬라이드 수(토큰 잘림 방지)


def plan_from_outline(generate_fn: Callable[[str, str, int], str],
                      outline_md: str, deck_title: str,
                      subtitle: str = "") -> List[Dict]:
    """개요를 슬라이드플랜으로 변환. 개요를 청크로 나눠 변환해 개수 잘림을 막는다.

    표지(cover)는 LLM이 아니라 여기서 맨 앞에 자동 추가한다. 청크가 실패하면
    그 청크만 블록 파싱(bullets)으로 폴백 → 전체 슬라이드 수는 항상 보존.
    """
    blocks = _outline_blocks(outline_md)
    if not blocks:
        return _fallback_plan(outline_md, deck_title)

    chunks = [blocks[i:i + _CHUNK] for i in range(0, len(blocks), _CHUNK)]
    out: List[Dict] = []
    for ci, ch in enumerate(chunks):
        chunk_md = "\n\n".join(ch)
        user = (f"{_ART_RULES}\n\n[덱 제목]\n{deck_title}\n"
                f"[부분 {ci + 1}/{len(chunks)} · 이 부분의 슬라이드 {len(ch)}개를 빠짐없이 변환]\n\n"
                f"[슬라이드 개요]\n{chunk_md}")
        print(f"[art] 청크 {ci + 1}/{len(chunks)} 생성 중… ({len(ch)}장 · gpt 응답 대기, ~1–2분)", flush=True)
        try:
            raw = generate_fn(_ART_SYS, user, 12000)
        except Exception as e:  # noqa: BLE001
            print(f"[art] 청크 {ci + 1} 생성 오류: {e}", flush=True)
            raw = ""
        arr = _extract_json_array(raw)
        got = [x for x in (_normalize_slide(s) for s in arr) if x] if isinstance(arr, list) else []
        # 개수가 모자라면 부족분을 블록 파싱으로 보충(개수 보존)
        if len(got) < len(ch):
            print(f"[art] 청크 {ci + 1}: {len(got)}/{len(ch)} → 부족분 블록 폴백", flush=True)
            got += [_block_to_bullets(b) for b in ch[len(got):]]
        out += got[:len(ch)]

    out.insert(0, {"type": "cover", "title": deck_title, "chip": subtitle})
    return out


def image_queries(plan: List[Dict]) -> Dict[int, str]:
    """photo 타입 슬라이드의 {인덱스: 영어 검색어}."""
    out = {}
    for i, s in enumerate(plan):
        if (s.get("type") == "photo") and s.get("image_query"):
            out[i] = s["image_query"]
    return out
